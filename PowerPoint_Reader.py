import glob
from pptx import Presentation
import pandas as pd

text_runs=[]


prs = Presentation('O:/test/IDW 1.0.pptx')


#prs = Presentation((path_to_presentation))
# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        tbl = shape.table
        row_count = len(tbl.rows)
        col_count = len(tbl.columns)
        try:
            cell = tbl.cell(0,1)
            print(cell.text_frame.parapgraphs)
            if 'Key Milestones' in cell.text_frame.paragraphs:
                for r in range(0, row_count):
                    for c in range(0, col_count):
                        cell = tbl.cell(r,c)
                        paragraphs = cell.text_frame.paragraphs
                        for paragraph in paragraphs:
                            for run in paragraph.runs:
                                print(f'{run.text}, col {c}, row, {r}')
                                text_runs.append(run.text)
        except:
            print('no')

'''
print('start')
for slide in prs.slides:
    #print('1')
    for shape in slide.shapes:
        #indicator == 0
        #print('2')
        if not shape.has_table:
            continue
        tbl=shape.table
        #cell=tbl.cell(0,0)
        paragraphs=cell.text_frame.paragraphs
        #for paragraph in paragraphs:
        #    for run in paragraph.runs:
        #        print(run.text)
            #    if "Key Milestones / Deliverables" in run.text:
            #        print('yes')
        row_count=len(tbl.rows)
        col_count=len(tbl.columns)
        for r in range(0,row_count):
            for c in range(0, col_count):
                cell = tbl.cell(r,c)
                parapgraphs=cell.text_frame.parapgraphs
                for paragraph in paragraphs:
                    for run in paragraph.runs:
                        prinit(f'{run.text}, col {c}, row, {r}')
                        text_runs.append(run.text)
                    #else:
                    #    print('no')
'''
